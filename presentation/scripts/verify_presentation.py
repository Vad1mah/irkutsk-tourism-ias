"""Verify built presentation: slide count, titles, content preview, overflow checks."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

OUTPUT = Path(__file__).resolve().parents[1] / "Защита_ВКР_Исполатов.pptx"


def main() -> None:
    prs = Presentation(str(OUTPUT))
    sw_cm = Emu(prs.slide_width).cm
    sh_cm = Emu(prs.slide_height).cm
    print(f"File: {OUTPUT}")
    print(f"Slide size: {sw_cm:.2f} × {sh_cm:.2f} cm")
    print(f"Slide count: {len(prs.slides)}")
    print()

    overflow_warnings: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        print(f"━━ Slide {idx} (layout: {slide.slide_layout.name!r}) ━━━━━━━━━━━━━━━━━━")
        text_pieces: list[str] = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    t = r.text.strip()
                    if t:
                        text_pieces.append(t)
            # Overflow check
            try:
                x = Emu(sh.left).cm
                y = Emu(sh.top).cm
                w = Emu(sh.width).cm
                h = Emu(sh.height).cm
                if x + w > sw_cm + 0.1:
                    overflow_warnings.append(
                        f"  Slide {idx}: shape {sh.name!r} right edge at {x + w:.2f} cm > slide width {sw_cm:.2f}"
                    )
                if y + h > sh_cm + 0.1:
                    overflow_warnings.append(
                        f"  Slide {idx}: shape {sh.name!r} bottom edge at {y + h:.2f} cm > slide height {sh_cm:.2f}"
                    )
                if x < -0.1 or y < -0.1:
                    overflow_warnings.append(
                        f"  Slide {idx}: shape {sh.name!r} negative position x={x:.2f} y={y:.2f}"
                    )
            except (TypeError, AttributeError):
                pass

        full_text = " | ".join(text_pieces)
        # Trim for display
        display = full_text if len(full_text) < 200 else full_text[:200] + " ..."
        print(f"  Text: {display}")
        print(f"  Shape count: {len(slide.shapes)}")
        print()

    print("━━ Overflow checks ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    if overflow_warnings:
        for w in overflow_warnings:
            print(w)
    else:
        print("  No overflow detected — all shapes fit within slide bounds.")


if __name__ == "__main__":
    main()
