"""Inspect IGU template to understand its structure: slide size, layouts, placeholders, fonts."""
from pptx import Presentation
from pptx.util import Emu

TEMPLATE = r"C:\Users\Admin\Desktop\Diplom\docs\vkr\Шаблон презентации ИГУ.pptx"


def main() -> None:
    prs = Presentation(TEMPLATE)

    print("=" * 70)
    print("SLIDE SIZE")
    print("=" * 70)
    print(f"  width  = {prs.slide_width} EMU = {Emu(prs.slide_width).inches:.2f}\" = {Emu(prs.slide_width).cm:.2f} cm")
    print(f"  height = {prs.slide_height} EMU = {Emu(prs.slide_height).inches:.2f}\" = {Emu(prs.slide_height).cm:.2f} cm")
    if prs.slide_width / prs.slide_height > 1.5:
        print("  -> 16:9 aspect ratio")
    else:
        print("  -> 4:3 aspect ratio")

    print()
    print("=" * 70)
    print(f"SLIDE LAYOUTS ({len(prs.slide_layouts)} total)")
    print("=" * 70)
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"\n[{idx}] Layout name: {layout.name!r}")
        for ph in layout.placeholders:
            print(f"    placeholder idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, name={ph.name!r}")
            if ph.has_text_frame:
                for p_idx, p in enumerate(ph.text_frame.paragraphs):
                    for r_idx, r in enumerate(p.runs):
                        font = r.font
                        size = font.size.pt if font.size else "inherit"
                        print(f"        p{p_idx}.r{r_idx}: {r.text!r} | font.name={font.name} size={size} bold={font.bold}")

    print()
    print("=" * 70)
    print(f"EXISTING SLIDES IN TEMPLATE ({len(prs.slides)} total)")
    print("=" * 70)
    for idx, slide in enumerate(prs.slides):
        print(f"\n[{idx}] uses layout: {slide.slide_layout.name!r}")
        for sh in slide.shapes:
            kind = sh.shape_type
            name = sh.name
            if sh.has_text_frame:
                texts = [r.text for p in sh.text_frame.paragraphs for r in p.runs]
                preview = " | ".join(t for t in texts if t.strip())[:80]
                print(f"    {kind} name={name!r} text={preview!r}")
            else:
                print(f"    {kind} name={name!r}")

    print()
    print("=" * 70)
    print("SLIDE MASTERS")
    print("=" * 70)
    for m_idx, master in enumerate(prs.slide_masters):
        print(f"\n[{m_idx}] master")
        for sh in master.shapes:
            if sh.has_text_frame:
                texts = [r.text for p in sh.text_frame.paragraphs for r in p.runs]
                preview = " | ".join(t for t in texts if t.strip())[:80]
                if preview:
                    print(f"    shape={sh.name!r} text={preview!r}")


if __name__ == "__main__":
    main()
