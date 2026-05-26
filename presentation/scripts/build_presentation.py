"""Build defense presentation from IGU template.

Run:
    PYTHONIOENCODING=utf-8 backend/venv/Scripts/python.exe presentation/scripts/build_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

# -------------------------------------------------------------------------
# Paths
# -------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = ROOT / "docs" / "vkr" / "Шаблон презентации ИГУ.pptx"
OUTPUT = ROOT / "presentation" / "Защита_ВКР_Исполатов.pptx"

# -------------------------------------------------------------------------
# Palette
# -------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
GREEN = RGBColor(0x2D, 0x8C, 0x5F)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

FONT = "Calibri"


# -------------------------------------------------------------------------
# Helpers
# -------------------------------------------------------------------------
def remove_all_slides(prs: Presentation) -> None:
    """Remove every slide from the presentation, keeping masters and layouts."""
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    rIds = [sldId.rId for sldId in list(sldIdLst)]
    for rId in rIds:
        prs.part.drop_rel(rId)
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)


def remove_placeholders(slide) -> None:
    """Remove all placeholders so we can draw our own composition."""
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def add_textbox(
    slide,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = FONT,
    letter_spacing: int | None = None,
):
    """Add a textbox with single paragraph styling."""
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(letter_spacing))
    return tb


def add_multiline_textbox(
    slide,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    lines: list[tuple[str, dict]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    """Add a textbox with several paragraphs, each line is (text, style_kwargs)."""
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    tf.vertical_anchor = anchor
    for i, (text, style) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = style.get("align", align)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if style.get("space_before") is not None:
            p.space_before = Pt(style["space_before"])
        if style.get("space_after") is not None:
            p.space_after = Pt(style["space_after"])
        run = p.add_run()
        run.text = text
        run.font.name = style.get("font_name", FONT)
        run.font.size = Pt(style.get("font_size", 18))
        run.font.bold = style.get("bold", False)
        run.font.color.rgb = style.get("color", BLACK)
        if style.get("letter_spacing") is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(style["letter_spacing"]))
    return tb


def add_rect(
    slide,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    *,
    fill: RGBColor | None = WHITE,
    line: RGBColor | None = DARK_BLUE,
    line_width: float = 1.0,
    shape_type=MSO_SHAPE.RECTANGLE,
    text: str | None = None,
    font_size: int = 18,
    bold: bool = False,
    text_color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.CENTER,
):
    """Add a rectangle (or other autoshape) with optional centered text."""
    shape = slide.shapes.add_shape(
        shape_type, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    if text is not None:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        tf.margin_top = Cm(0.1)
        tf.margin_bottom = Cm(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape


def add_rounded_rect(slide, *args, **kwargs):
    """Convenience: rounded rectangle."""
    kwargs.setdefault("shape_type", MSO_SHAPE.ROUNDED_RECTANGLE)
    return add_rect(slide, *args, **kwargs)


def add_down_arrow(slide, center_cm: float, top_cm: float, height_cm: float = 0.8) -> None:
    """Draw a small down-pointing arrow symbol."""
    width = 0.6
    add_textbox(
        slide,
        center_cm - width / 2,
        top_cm,
        width,
        height_cm,
        "▼",
        font_size=20,
        color=DARK_BLUE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_footer_pageno(slide, page_num: int, total: int = 11) -> None:
    """Page indicator in the right-bottom corner."""
    add_textbox(
        slide,
        31.5,
        18.2,
        2.0,
        0.6,
        f"{page_num} / {total}",
        font_size=10,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


# -------------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------------
def build_slide_01_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Титульный слайд
    remove_placeholders(slide)

    # Top: Институт
    add_textbox(
        slide, 0, 1.2, 33.87, 0.9,
        "ИРКУТСКИЙ ГОСУДАРСТВЕННЫЙ УНИВЕРСИТЕТ",
        font_size=14, color=GRAY, align=PP_ALIGN.CENTER, letter_spacing=200,
    )
    add_textbox(
        slide, 0, 2.05, 33.87, 0.7,
        "Факультет бизнес-коммуникаций и информатики",
        font_size=13, color=GRAY, align=PP_ALIGN.CENTER,
    )

    # Section label
    add_textbox(
        slide, 0, 4.0, 33.87, 0.7,
        "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА",
        font_size=16, color=DARK_BLUE, align=PP_ALIGN.CENTER, letter_spacing=300, bold=True,
    )
    # Divider line
    add_rect(slide, 13.0, 4.85, 7.87, 0.05, fill=DARK_BLUE, line=None)

    # Theme
    theme = (
        "Интеллектуальная система прогнозирования заполняемости средств размещения "
        "на основе данных туристических агрегаторов и событийной активности в регионе"
    )
    add_textbox(
        slide, 2.0, 5.5, 29.87, 4.0, theme,
        font_size=24, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Author / supervisor block
    add_multiline_textbox(
        slide, 9.5, 12.5, 14.87, 3.5,
        [
            ("Выполнил:", {"font_size": 14, "color": GRAY, "letter_spacing": 150, "space_after": 4}),
            ("Исполатов Вадим Павлович, группа 14322-ДБ",
             {"font_size": 18, "bold": True, "color": BLACK, "space_after": 14}),
            ("Научный руководитель:", {"font_size": 14, "color": GRAY, "letter_spacing": 150, "space_after": 4}),
            ("ст. преп. Пестова Ю. В.",
             {"font_size": 18, "bold": True, "color": BLACK}),
        ],
        align=PP_ALIGN.CENTER,
    )

    # Footer
    add_textbox(
        slide, 0, 17.4, 33.87, 0.8,
        "Иркутск, 2026", font_size=14, color=GRAY, align=PP_ALIGN.CENTER, letter_spacing=200,
    )


def build_slide_02_goals(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Цели и задачи")
    add_footer_pageno(slide, 2)

    # Цель block
    add_textbox(
        slide, 1.5, 2.7, 5.0, 0.7, "ЦЕЛЬ",
        font_size=16, bold=True, color=DARK_BLUE, letter_spacing=300,
    )
    goal_text = (
        "Проектирование и разработка интеллектуальной аналитической системы "
        "прогнозирования заполняемости средств размещения с интерфейсом визуальной "
        "аналитики на основе данных туристических агрегаторов и событийной активности "
        "в регионе"
    )
    add_rect(
        slide, 1.5, 3.5, 30.87, 2.7,
        fill=LIGHT_GRAY, line=None, text=goal_text,
        font_size=15, text_color=BLACK, align=PP_ALIGN.LEFT,
    )

    # Задачи block
    add_textbox(
        slide, 1.5, 6.7, 5.0, 0.7, "ЗАДАЧИ",
        font_size=16, bold=True, color=DARK_BLUE, letter_spacing=300,
    )

    tasks = [
        "Исследовать инструменты мониторинга и систем управления доходностью (RMS).",
        "Обосновать выбор технологий и реализовать архитектуру серверной части.",
        "Спроектировать подсистему автоматизированного сбора данных.",
        "Разработать ансамблевый метод прогнозирования с расчётом RMS-метрик.",
        "Спроектировать архитектуру ИАС с интеллектуальным агентом и веб-интерфейсом для трёх профильных сегментов.",
    ]
    lines = []
    for i, t in enumerate(tasks, start=1):
        lines.append(
            (f"{i}.  {t}", {
                "font_size": 15,
                "color": BLACK,
                "space_after": 6,
            })
        )
    add_multiline_textbox(slide, 1.5, 7.5, 30.87, 7.5, lines, line_spacing=1.15)

    # Якорь-плашка с цифрами
    add_rect(slide, 0, 16.3, 33.87, 1.5, fill=DARK_BLUE, line=None)
    add_textbox(
        slide, 0, 16.3, 33.87, 1.5,
        "1 428 объектов   ·   450+ событий   ·   RMSE 2,67 п. п.   ·   12 RMS-инструментов",
        font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def build_slide_03_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Архитектура решения")
    add_footer_pageno(slide, 3)

    # Top: sources row label
    add_textbox(
        slide, 1.5, 2.8, 10.0, 0.6, "ИСТОЧНИКИ",
        font_size=13, color=GRAY, letter_spacing=300, bold=True,
    )

    # Three source boxes
    source_y = 3.5
    source_w = 6.5
    source_h = 1.5
    source_gap = 1.5
    total_w = 3 * source_w + 2 * source_gap  # 22.5
    start_x = (33.87 - total_w) / 2  # 5.685
    for i, name in enumerate(["События", "Отели", "Погода"]):
        x = start_x + i * (source_w + source_gap)
        add_rect(slide, x, source_y, source_w, source_h,
                 fill=WHITE, line=GRAY, line_width=1.0,
                 text=name, font_size=16, text_color=BLACK, bold=False)

    # Three arrows down to "Сбор данных"
    arrow_top = source_y + source_h + 0.1
    for i in range(3):
        x = start_x + i * (source_w + source_gap) + source_w / 2
        add_down_arrow(slide, x, arrow_top, 0.8)

    # Layer 1: Сбор данных
    layer_x = 4.0
    layer_w = 25.87
    layer_h = 1.7
    sbor_y = source_y + source_h + 1.0
    add_rounded_rect(
        slide, layer_x, sbor_y, layer_w, layer_h,
        fill=WHITE, line=DARK_BLUE, line_width=2.0,
    )
    add_multiline_textbox(
        slide, layer_x, sbor_y, layer_w, layer_h,
        [
            ("СБОР ДАННЫХ", {"font_size": 20, "bold": True, "color": DARK_BLUE,
                             "align": PP_ALIGN.CENTER, "letter_spacing": 200, "space_after": 4}),
            ("автоматический,  24 / 7", {"font_size": 14, "color": GRAY,
                                          "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Arrow down
    add_down_arrow(slide, 33.87 / 2, sbor_y + layer_h + 0.1, 0.7)

    # Layer 2: Аналитическое ядро
    yadro_y = sbor_y + layer_h + 0.95
    add_rounded_rect(
        slide, layer_x, yadro_y, layer_w, layer_h,
        fill=WHITE, line=DARK_BLUE, line_width=2.0,
    )
    add_multiline_textbox(
        slide, layer_x, yadro_y, layer_w, layer_h,
        [
            ("АНАЛИТИЧЕСКОЕ ЯДРО", {"font_size": 20, "bold": True, "color": DARK_BLUE,
                                    "align": PP_ALIGN.CENTER, "letter_spacing": 200, "space_after": 4}),
            ("ML-прогноз  ·  RMS-метрики  ·  AI-объяснение",
             {"font_size": 14, "color": GRAY, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    add_down_arrow(slide, 33.87 / 2, yadro_y + layer_h + 0.1, 0.7)

    # Layer 3: Интерфейсы доступа
    ui_y = yadro_y + layer_h + 0.95
    add_rounded_rect(
        slide, layer_x, ui_y, layer_w, layer_h,
        fill=WHITE, line=DARK_BLUE, line_width=2.0,
    )
    add_multiline_textbox(
        slide, layer_x, ui_y, layer_w, layer_h,
        [
            ("ИНТЕРФЕЙСЫ ДОСТУПА", {"font_size": 20, "bold": True, "color": DARK_BLUE,
                                    "align": PP_ALIGN.CENTER, "letter_spacing": 200, "space_after": 4}),
            ("Дашборд  ·  AI-чат",
             {"font_size": 14, "color": GRAY, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    add_down_arrow(slide, 33.87 / 2, ui_y + layer_h + 0.1, 0.7)

    # Bottom: pользователи row
    users_y = ui_y + layer_h + 0.95
    add_textbox(
        slide, 1.5, users_y - 0.6, 10.0, 0.5, "ПОЛЬЗОВАТЕЛИ",
        font_size=13, color=GRAY, letter_spacing=300, bold=True,
    )
    users = ["Отельер", "Администрация", "Исследователь"]
    user_w = 6.5
    user_h = 1.3
    user_gap = 1.5
    user_total = 3 * user_w + 2 * user_gap
    user_start = (33.87 - user_total) / 2
    for i, name in enumerate(users):
        x = user_start + i * (user_w + user_gap)
        add_rect(slide, x, users_y, user_w, user_h,
                 fill=DARK_BLUE, line=None,
                 text=name, font_size=16, bold=True, text_color=WHITE)


def build_slide_04_research(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 1.  Исследование предметной области")
    add_footer_pageno(slide, 4)

    rows = [
        ("Направление", "Выбор и почему"),
        ("RMS-метрики", "Свой расчёт на открытых данных — коммерческие системы регион не покрывают"),
        ("Метод прогноза", "Ансамбль 3 ML-моделей — ARIMA не принимает экзогенные регрессоры"),
        ("LLM-фреймворк", "LangGraph — граф состояний предсказуемее «чёрного ящика» LangChain Agents"),
        ("Источники", "11 потоков, унифицированный сборщик"),
    ]
    table_top = 3.5
    col1_w = 7.5
    col2_w = 22.0
    table_left = (33.87 - col1_w - col2_w) / 2

    # header row
    header_h = 1.0
    add_rect(slide, table_left, table_top, col1_w, header_h,
             fill=DARK_BLUE, line=None,
             text=rows[0][0], font_size=16, bold=True, text_color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, table_left + col1_w, table_top, col2_w, header_h,
             fill=DARK_BLUE, line=None,
             text=rows[0][1], font_size=16, bold=True, text_color=WHITE, align=PP_ALIGN.LEFT)

    # body rows
    row_h = 1.9
    y = table_top + header_h
    for i, (col1, col2) in enumerate(rows[1:]):
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, table_left, y, col1_w, row_h,
                 fill=bg, line=GRAY, line_width=0.5,
                 text=col1, font_size=15, text_color=GRAY, align=PP_ALIGN.LEFT, bold=False)
        add_rect(slide, table_left + col1_w, y, col2_w, row_h,
                 fill=bg, line=GRAY, line_width=0.5,
                 text=col2, font_size=16, bold=True, text_color=DARK_BLUE, align=PP_ALIGN.LEFT)
        # left padding fix: re-add text with padding via inner textbox? skipping — anchor centers vertically
        y += row_h


def build_slide_05_tech(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 2.  Технологический стек")
    add_footer_pageno(slide, 5)

    # 2x2 grid of stack groups
    grid_x = 1.5
    grid_y = 3.0
    cell_w = 15.0
    cell_h = 5.4
    gap_x = 0.87
    gap_y = 0.6

    groups = [
        # (col, row, label, items, accent)
        ("FRONTEND", [
            "React 18  +  TypeScript",
            "Vite,  Tailwind CSS",
            "Recharts  (графики)",
            "Yandex Maps  (карта)",
        ]),
        ("BACKEND", [
            "Python 3.11",
            "FastAPI  (async REST)",
            "SQLAlchemy 2.0  +  asyncpg",
            "Pydantic v2",
        ]),
        ("ХРАНЕНИЕ", [
            "PostgreSQL 16  (основная БД)",
            "Redis 7  (кэш, rate-limit)",
            "ChromaDB  (векторная база, RAG)",
        ]),
        ("ИНФРАСТРУКТУРА", [
            "Docker  (multi-stage сборка)",
            "APScheduler  (планировщик)",
            "Alembic  (каркас миграций)",
        ]),
    ]

    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for (label, items), (col, row) in zip(groups, positions):
        x = grid_x + col * (cell_w + gap_x)
        y = grid_y + row * (cell_h + gap_y)
        add_rounded_rect(slide, x, y, cell_w, cell_h,
                         fill=WHITE, line=DARK_BLUE, line_width=1.5)
        # Header strip
        add_rect(slide, x, y, cell_w, 1.0,
                 fill=DARK_BLUE, line=None,
                 text=label, font_size=14, bold=True, text_color=WHITE,
                 align=PP_ALIGN.CENTER)
        # Items
        item_lines = [
            (item, {"font_size": 15, "color": BLACK, "space_after": 6})
            for item in items
        ]
        add_multiline_textbox(
            slide, x + 0.6, y + 1.2, cell_w - 1.2, cell_h - 1.4,
            item_lines, line_spacing=1.15,
        )

    # Внизу: тонкая плашка про миграцию (как факт, не как акцент слайда)
    plate_y = grid_y + 2 * cell_h + gap_y + 0.6
    plate_w = 30.87
    plate_x = (33.87 - plate_w) / 2
    plate_h = 1.6
    add_rounded_rect(slide, plate_x, plate_y, plate_w, plate_h,
                     fill=LIGHT_GRAY, line=None)
    add_multiline_textbox(
        slide, plate_x + 0.6, plate_y + 0.2, plate_w - 1.2, plate_h - 0.4,
        [
            ("ОТДЕЛЬНО ВЫПОЛНЕНО:  миграция БД  Yandex YDB  →  PostgreSQL 16",
             {"font_size": 13, "bold": True, "color": DARK_BLUE,
              "letter_spacing": 200, "align": PP_ALIGN.CENTER, "space_after": 3}),
            ("локальная разработка через Docker  ·  асинхронный ORM  ·  независимость от облачного провайдера",
             {"font_size": 12, "color": GRAY, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )


def build_slide_06_data(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 3.  Подсистема сбора данных")
    add_footer_pageno(slide, 6)

    # Three big number plates
    plate_y = 4.5
    plate_w = 9.0
    plate_h = 5.5
    plate_gap = 1.0
    total = 3 * plate_w + 2 * plate_gap
    start_x = (33.87 - total) / 2

    big_numbers = [
        ("1 428", "объектов размещения"),
        ("450+", "событий региона"),
        ("11", "внешних источников"),
    ]
    for i, (num, label) in enumerate(big_numbers):
        x = start_x + i * (plate_w + plate_gap)
        add_rounded_rect(slide, x, plate_y, plate_w, plate_h,
                         fill=WHITE, line=DARK_BLUE, line_width=1.5)
        add_textbox(slide, x, plate_y + 0.5, plate_w, 3.4,
                    num, font_size=66, bold=True, color=GREEN,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, plate_y + 4.0, plate_w, 1.2,
                    label, font_size=16, color=GRAY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Automation plate
    auto_y = plate_y + plate_h + 1.0
    auto_h = 2.4
    auto_w = 28.0
    auto_x = (33.87 - auto_w) / 2
    add_rounded_rect(slide, auto_x, auto_y, auto_w, auto_h,
                     fill=LIGHT_GRAY, line=None)
    add_multiline_textbox(
        slide, auto_x + 0.5, auto_y + 0.3, auto_w - 1.0, auto_h - 0.6,
        [
            ("АВТОМАТИЧЕСКИЙ СБОР ПО РАСПИСАНИЮ",
             {"font_size": 13, "bold": True, "color": DARK_BLUE,
              "letter_spacing": 300, "align": PP_ALIGN.CENTER, "space_after": 6}),
            ("События каждые 6 ч  ·  Отели каждые 2 ч  ·  Погода каждые 3 ч  ·  Telegram каждый час",
             {"font_size": 16, "color": BLACK, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Bottom small note
    add_textbox(
        slide, 1.5, auto_y + auto_h + 0.5, 30.87, 0.7,
        "Унифицированный сборщик из 17 файлов парсеров  ·  LLM-фильтрация Telegram (ministral-8b)",
        font_size=13, color=GRAY, align=PP_ALIGN.CENTER,
    )


def build_slide_07_ensemble(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 4.  Ансамблевый метод прогнозирования")
    add_footer_pageno(slide, 7)

    # ВХОД block
    input_y = 2.7
    input_h = 1.6
    input_w = 26.0
    input_x = (33.87 - input_w) / 2
    add_rounded_rect(slide, input_x, input_y, input_w, input_h,
                     fill=WHITE, line=DARK_BLUE, line_width=1.5)
    add_multiline_textbox(
        slide, input_x, input_y, input_w, input_h,
        [
            ("ВХОД:  38 признаков  (6 групп)",
             {"font_size": 16, "bold": True, "color": DARK_BLUE,
              "align": PP_ALIGN.CENTER, "space_after": 2}),
            ("календарные  ·  лаги  ·  скользящие  ·  погода  ·  события  ·  тренд",
             {"font_size": 13, "color": GRAY, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Arrows down to 3 models
    arrow_y = input_y + input_h + 0.15
    add_down_arrow(slide, 33.87 / 2, arrow_y, 0.5)

    # Three models
    models = [
        ("Prophet", "сезонность\n+ долгосрочный тренд"),
        ("NeuralProphet", "нелинейные паттерны\nчерез автогрегрессию"),
        ("XGBoost", "градиентный бустинг,\nсобытия и погода"),
    ]
    m_y = arrow_y + 0.8
    m_w = 8.5
    m_h = 3.0
    m_gap = 1.0
    m_total = 3 * m_w + 2 * m_gap
    m_start = (33.87 - m_total) / 2
    for i, (name, desc) in enumerate(models):
        x = m_start + i * (m_w + m_gap)
        add_rounded_rect(slide, x, m_y, m_w, m_h,
                         fill=WHITE, line=DARK_BLUE, line_width=1.5)
        add_textbox(slide, x, m_y + 0.3, m_w, 0.9,
                    name, font_size=22, bold=True, color=DARK_BLUE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, m_y + 1.3, m_w, 1.5,
                    desc, font_size=13, color=GRAY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Arrows down to calibration block
    calib_arrow_y = m_y + m_h + 0.15
    add_down_arrow(slide, 33.87 / 2, calib_arrow_y, 0.5)

    # Calibration block - highlighted
    calib_y = calib_arrow_y + 0.7
    calib_h = 1.8
    calib_w = 26.0
    calib_x = (33.87 - calib_w) / 2
    add_rounded_rect(slide, calib_x, calib_y, calib_w, calib_h,
                     fill=GREEN_LIGHT, line=GREEN, line_width=2.5)
    add_multiline_textbox(
        slide, calib_x, calib_y, calib_w, calib_h,
        [
            ("АДАПТИВНАЯ КАЛИБРОВКА",
             {"font_size": 17, "bold": True, "color": DARK_BLUE,
              "letter_spacing": 300, "align": PP_ALIGN.CENTER, "space_after": 4}),
            ("веса  ~  1 / RMSE  на 14 днях  ·  пересчёт раз в час",
             {"font_size": 14, "color": BLACK, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Arrow down to output
    out_arrow_y = calib_y + calib_h + 0.15
    add_down_arrow(slide, 33.87 / 2, out_arrow_y, 0.5)

    # Output
    out_y = out_arrow_y + 0.7
    out_h = 1.2
    out_w = 26.0
    out_x = (33.87 - out_w) / 2
    add_rounded_rect(slide, out_x, out_y, out_w, out_h,
                     fill=DARK_BLUE, line=None)
    add_textbox(slide, out_x, out_y, out_w, out_h,
                "ВЫХОД:  прогноз  +  интервал доверия  80 %",
                font_size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide_08_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 4.  Метрики ансамбля")
    add_footer_pageno(slide, 8)

    # Subtitle
    add_textbox(
        slide, 0, 2.8, 33.87, 0.7,
        "Иркутский район, горизонт 7 дней",
        font_size=15, color=GRAY, align=PP_ALIGN.CENTER,
    )

    # Table
    rows = [
        ("Модель", "RMSE (п. п.)", None),
        ("ENSEMBLE", "2,67", "highlight"),
        ("Prophet", "5,58", None),
        ("NeuralProphet", "5,92", None),
        ("XGBoost", "6,11", None),
    ]
    t_left = 7.0
    t_w = 19.87
    col1_w = 12.0
    col2_w = t_w - col1_w
    row_h_header = 1.0
    row_h_body = 1.3
    y = 4.0
    # header
    add_rect(slide, t_left, y, col1_w, row_h_header,
             fill=DARK_BLUE, line=None,
             text=rows[0][0], font_size=15, bold=True, text_color=WHITE)
    add_rect(slide, t_left + col1_w, y, col2_w, row_h_header,
             fill=DARK_BLUE, line=None,
             text=rows[0][1], font_size=15, bold=True, text_color=WHITE)
    y += row_h_header

    for name, val, mark in rows[1:]:
        if mark == "highlight":
            bg = GREEN_LIGHT
            text_color = DARK_BLUE
            num_color = GREEN
            font_size_name = 22
            font_size_num = 28
            bold = True
        else:
            bg = WHITE
            text_color = BLACK
            num_color = BLACK
            font_size_name = 18
            font_size_num = 20
            bold = False
        add_rect(slide, t_left, y, col1_w, row_h_body,
                 fill=bg, line=GRAY, line_width=0.5,
                 text=name, font_size=font_size_name, bold=bold, text_color=text_color)
        add_rect(slide, t_left + col1_w, y, col2_w, row_h_body,
                 fill=bg, line=GRAY, line_width=0.5,
                 text=val, font_size=font_size_num, bold=True, text_color=num_color)
        y += row_h_body

    # Bottom plate: business meaning
    plate_y = y + 1.0
    plate_w = 28.0
    plate_h = 4.2
    plate_x = (33.87 - plate_w) / 2
    add_rounded_rect(slide, plate_x, plate_y, plate_w, plate_h,
                     fill=DARK_BLUE, line=None)
    add_multiline_textbox(
        slide, plate_x + 0.5, plate_y + 0.3, plate_w - 1.0, plate_h - 0.6,
        [
            ("ЧТО ЭТО ЗНАЧИТ ДЛЯ ОТЕЛЬЕРА",
             {"font_size": 13, "bold": True, "color": WHITE,
              "letter_spacing": 300, "align": PP_ALIGN.CENTER, "space_after": 8}),
            ("Прогноз 65 %  →  факт между 62 и 68 %",
             {"font_size": 20, "bold": True, "color": WHITE,
              "align": PP_ALIGN.CENTER, "space_after": 6}),
            ("Точность достаточна для решений по тарифной политике",
             {"font_size": 15, "color": WHITE, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )


def build_slide_09_agent(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 5.  Архитектура ИАС и LLM-агент")
    add_footer_pageno(slide, 9)

    # Left column: state graph
    left_x = 1.5
    left_w = 11.0
    col_y = 2.8
    col_h = 12.2
    add_rounded_rect(slide, left_x, col_y, left_w, col_h,
                     fill=LIGHT_GRAY, line=None)
    add_textbox(slide, left_x, col_y + 0.3, left_w, 0.6,
                "ГРАФ СОСТОЯНИЙ LANGGRAPH",
                font_size=13, bold=True, color=DARK_BLUE,
                letter_spacing=300, align=PP_ALIGN.CENTER)

    # State graph nodes
    node_w = 4.0
    node_h = 1.3
    cx = left_x + left_w / 2
    # START
    start_y = col_y + 1.5
    add_rect(slide, cx - node_w / 2, start_y, node_w, node_h,
             shape_type=MSO_SHAPE.OVAL,
             fill=GREEN, line=None, text="START",
             font_size=14, bold=True, text_color=WHITE)
    # arrow
    add_down_arrow(slide, cx, start_y + node_h + 0.05, 0.6)
    # model
    model_y = start_y + node_h + 0.7
    add_rect(slide, cx - node_w / 2, model_y, node_w, node_h,
             shape_type=MSO_SHAPE.OVAL,
             fill=DARK_BLUE, line=None, text="model",
             font_size=14, bold=True, text_color=WHITE)
    add_down_arrow(slide, cx, model_y + node_h + 0.05, 0.6)
    # tools
    tools_y = model_y + node_h + 0.7
    add_rect(slide, cx - node_w / 2, tools_y, node_w, node_h,
             shape_type=MSO_SHAPE.OVAL,
             fill=DARK_BLUE, line=None, text="tools",
             font_size=14, bold=True, text_color=WHITE)
    # circular loop label (text only)
    add_textbox(slide, cx + node_w / 2 + 0.15, tools_y - 0.3, 2.0, 1.8,
                "↺",
                font_size=36, color=DARK_BLUE, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_down_arrow(slide, cx, tools_y + node_h + 0.05, 0.6)
    # END
    end_y = tools_y + node_h + 0.7
    add_rect(slide, cx - node_w / 2, end_y, node_w, node_h,
             shape_type=MSO_SHAPE.OVAL,
             fill=RGBColor(0xC0, 0x3E, 0x3A), line=None, text="END",
             font_size=14, bold=True, text_color=WHITE)

    # Right column: dialogue scenario
    right_x = 13.5
    right_w = 18.87
    add_rounded_rect(slide, right_x, col_y, right_w, col_h,
                     fill=WHITE, line=DARK_BLUE, line_width=1.5)
    add_textbox(slide, right_x, col_y + 0.3, right_w, 0.6,
                "СЦЕНАРИЙ ИСПОЛЬЗОВАНИЯ",
                font_size=13, bold=True, color=DARK_BLUE,
                letter_spacing=300, align=PP_ALIGN.CENTER)

    dialog_top = col_y + 1.3
    # Otellier turn
    add_textbox(slide, right_x + 0.5, dialog_top, right_w - 1.0, 0.6,
                "ОТЕЛЬЕР:", font_size=12, bold=True, color=DARK_BLUE, letter_spacing=300)
    add_rounded_rect(slide, right_x + 0.5, dialog_top + 0.6, right_w - 1.0, 1.4,
                     fill=LIGHT_GRAY, line=None,
                     text="«Сравни Иркутский и Ольхонский за 30 дней»",
                     font_size=15, text_color=BLACK, align=PP_ALIGN.CENTER)
    # arrow
    add_down_arrow(slide, right_x + right_w / 2, dialog_top + 2.1, 0.5)
    # Agent turn
    agent_top = dialog_top + 2.8
    add_textbox(slide, right_x + 0.5, agent_top, right_w - 1.0, 0.6,
                "АГЕНТ:", font_size=12, bold=True, color=GRAY, letter_spacing=300)
    add_rounded_rect(slide, right_x + 0.5, agent_top + 0.6, right_w - 1.0, 1.4,
                     fill=LIGHT_GRAY, line=None,
                     text="вызывает compare_districts(),  get_revenue_metrics()",
                     font_size=13, text_color=BLACK, align=PP_ALIGN.CENTER)
    # arrow
    add_down_arrow(slide, right_x + right_w / 2, agent_top + 2.1, 0.5)
    # Answer
    ans_top = agent_top + 2.8
    add_textbox(slide, right_x + 0.5, ans_top, right_w - 1.0, 0.6,
                "ОТВЕТ:", font_size=12, bold=True, color=GREEN, letter_spacing=300)
    add_rounded_rect(slide, right_x + 0.5, ans_top + 0.6, right_w - 1.0, 2.7,
                     fill=GREEN_LIGHT, line=GREEN, line_width=1.0)
    add_multiline_textbox(
        slide, right_x + 0.5, ans_top + 0.6, right_w - 1.0, 2.7,
        [
            ("Ольхонский RevPAR  1 884 ₽,",
             {"font_size": 16, "bold": True, "color": DARK_BLUE,
              "align": PP_ALIGN.CENTER, "space_after": 3}),
            ("Иркутский  1 749 ₽",
             {"font_size": 16, "bold": True, "color": DARK_BLUE,
              "align": PP_ALIGN.CENTER, "space_after": 8}),
            ("Разница  135 ₽  с номера в день",
             {"font_size": 14, "color": BLACK, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Bottom strip
    strip_y = col_y + col_h + 0.4
    add_textbox(
        slide, 1.5, strip_y, 30.87, 0.6,
        "12 инструментов  ·  RAG ChromaDB  ·  Fallback: Groq → DeepSeek → Mistral  ·  Стресс-тест 19/19 OK",
        font_size=12, color=GRAY, align=PP_ALIGN.CENTER, bold=True,
    )


def build_slide_10_ui(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Этап 5.  Веб-интерфейс")
    add_footer_pageno(slide, 10)

    # Screenshot placeholder
    img_x = 3.5
    img_y = 2.8
    img_w = 26.87
    img_h = 10.5
    add_rect(slide, img_x, img_y, img_w, img_h,
             fill=LIGHT_GRAY, line=DARK_BLUE, line_width=1.5)
    add_multiline_textbox(
        slide, img_x, img_y, img_w, img_h,
        [
            ("[ СКРИНШОТ ТЕПЛОВОЙ КАРТЫ RMS-ДАШБОРДА ]",
             {"font_size": 16, "bold": True, "color": GRAY,
              "align": PP_ALIGN.CENTER, "space_after": 8}),
            ("дни недели × месяцы, штриховка на ячейках с малой выборкой",
             {"font_size": 13, "color": GRAY, "align": PP_ALIGN.CENTER, "space_after": 12}),
            ("→ заменить в день защиты свежим скриншотом из /analytics",
             {"font_size": 11, "color": DARK_BLUE, "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Bottom plate
    plate_y = img_y + img_h + 0.5
    plate_w = 28.0
    plate_h = 2.5
    plate_x = (33.87 - plate_w) / 2
    add_rounded_rect(slide, plate_x, plate_y, plate_w, plate_h,
                     fill=DARK_BLUE, line=None,
                     text="Тепловая карта показывает провалы загрузки по дням и сезонам —\nотельер видит, куда подкрутить цену",
                     font_size=16, bold=True, text_color=WHITE)

    # Bottom note
    add_textbox(
        slide, 1.5, plate_y + plate_h + 0.3, 30.87, 0.6,
        "8 страниц  ·  205 модульных тестов в 34 файлах  ·  9 E2E-сценариев  ·  persona walkthrough",
        font_size=12, color=GRAY, align=PP_ALIGN.CENTER,
    )


def build_slide_11_conclusion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    remove_placeholders(slide)
    add_slide_title(slide, "Заключение")
    add_footer_pageno(slide, 11)

    add_textbox(slide, 1.5, 2.7, 15.0, 0.6,
                "РЕЗУЛЬТАТЫ ПО ЗАДАЧАМ",
                font_size=14, bold=True, color=DARK_BLUE, letter_spacing=300)

    # Каждый пункт — вывод/результат, не пересказ процесса.
    # Структура: жирная мысль-результат + пояснение
    results = [
        ("Обоснован выбор технологического стека",
         "для работы с разнородными источниками данных и экзогенными факторами спроса"),
        ("Спроектирована масштабируемая серверная архитектура",
         "с асинхронной обработкой запросов и контейнерным развёртыванием стека"),
        ("Создан единый источник рыночной информации по региону:",
         "1 428 объектов размещения и 450+ событий из 11 источников"),
        ("Разработан ансамблевый метод прогнозирования",
         "с точностью RMSE 2,67 п. п. — превосходит каждую автономную модель"),
        ("Реализована B2B-платформа со встроенным интеллектуальным агентом",
         "для отельеров, региональной администрации и исследователей"),
    ]
    lines: list[tuple[str, dict]] = []
    for i, (head, body) in enumerate(results, start=1):
        lines.append(
            (f"✓   {i}.   {head}", {
                "font_size": 14, "bold": True, "color": DARK_BLUE,
                "space_after": 1,
            })
        )
        lines.append(
            (f"         {body}", {
                "font_size": 13, "color": BLACK,
                "space_after": 8,
            })
        )
    add_multiline_textbox(slide, 1.5, 3.5, 30.87, 8.5, lines, line_spacing=1.15)

    # Value plate
    plate_y = 12.0
    plate_w = 30.0
    plate_h = 4.3
    plate_x = (33.87 - plate_w) / 2
    add_rounded_rect(slide, plate_x, plate_y, plate_w, plate_h,
                     fill=WHITE, line=DARK_BLUE, line_width=2.5)
    add_textbox(slide, plate_x, plate_y + 0.3, plate_w, 0.8,
                "ГОТОВА К ПИЛОТУ",
                font_size=18, bold=True, color=DARK_BLUE,
                letter_spacing=400, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # divider line
    add_rect(slide, plate_x + 11.5, plate_y + 1.2, 7.0, 0.04,
             fill=DARK_BLUE, line=None)

    add_multiline_textbox(
        slide, plate_x + 1.0, plate_y + 1.5, plate_w - 2.0, plate_h - 1.8,
        [
            ("—  Отельеру:  RMS-метрики без коммерческой подписки 5 тыс ₽/мес",
             {"font_size": 13, "color": BLACK, "space_after": 4}),
            ("—  Администрации:  картина рынка без полугодовой задержки Росстата",
             {"font_size": 13, "color": BLACK, "space_after": 4}),
            ("—  Исследователю:  CSV-экспорт и прозрачная методология",
             {"font_size": 13, "color": BLACK}),
        ],
        anchor=MSO_ANCHOR.TOP,
    )

    # Bottom thank you
    add_textbox(
        slide, 0, 16.8, 33.87, 0.8,
        "СПАСИБО ЗА ВНИМАНИЕ.   Готов ответить на вопросы.",
        font_size=17, bold=True, color=DARK_BLUE,
        letter_spacing=200, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# -------------------------------------------------------------------------
# Title helper (used everywhere except slide 1)
# -------------------------------------------------------------------------
def add_slide_title(slide, text: str) -> None:
    add_textbox(
        slide, 1.5, 0.8, 30.87, 1.4, text,
        font_size=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT,
    )
    # accent underline
    add_rect(slide, 1.5, 2.15, 4.0, 0.08, fill=GREEN, line=None)


# -------------------------------------------------------------------------
# Main
# -------------------------------------------------------------------------
def main() -> None:
    print(f"Template: {TEMPLATE}")
    print(f"Output:   {OUTPUT}")

    prs = Presentation(str(TEMPLATE))
    print(f"Loaded template with {len(prs.slides)} sample slides.")

    remove_all_slides(prs)
    print(f"Cleared. Slides now: {len(prs.slides)}")

    print("Building 11 slides...")
    build_slide_01_title(prs)
    build_slide_02_goals(prs)
    build_slide_03_architecture(prs)
    build_slide_04_research(prs)
    build_slide_05_tech(prs)
    build_slide_06_data(prs)
    build_slide_07_ensemble(prs)
    build_slide_08_metrics(prs)
    build_slide_09_agent(prs)
    build_slide_10_ui(prs)
    build_slide_11_conclusion(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Final slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
